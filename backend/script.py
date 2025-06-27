
# import os
# import shutil
# from PyPDF2 import PdfReader
# from pptx import Presentation
# import docx
# from pathlib import Path

# def process_files(input_dir="old_docs", output_dir="docs"):
#     """Process all supported files and save as text"""
#     Path(output_dir).mkdir(exist_ok=True)
    
#     for filename in os.listdir(input_dir):
#         filepath = os.path.join(input_dir, filename)
#         output_path = os.path.join(output_dir, f"{Path(filename).stem}.txt")
        
#         try:
#             if filename.lower().endswith('.pdf'):
#                 text = extract_pdf_text(filepath)
#             elif filename.lower().endswith('.pptx'):
#                 text = extract_pptx_text(filepath)
#             elif filename.lower().endswith('.docx'):
#                 text = extract_docx_text(filepath)
#             elif filename.lower().endswith('.txt'):
#                 shutil.copy2(filepath, output_path)
#                 print(f"Copied: {filename} → {output_path}")
#                 continue
#             else:
#                 continue
            
#             with open(output_path, 'w', encoding='utf-8') as f:
#                 f.write(text)
#             print(f"Processed: {filename} → {output_path}")
            
#         except Exception as e:
#             print(f"Failed {filename}: {str(e)}")

# def extract_pdf_text(path):
#     with open(path, 'rb') as f:
#         return "\n".join(
#             p.extract_text() 
#             for p in PdfReader(f).pages 
#             if p.extract_text()
#         )

# def extract_pptx_text(path):
#     prs = Presentation(path)
#     return "\n".join(
#         shape.text.strip() 
#         for slide in prs.slides 
#         for shape in slide.shapes 
#         if hasattr(shape, "text") and shape.text.strip()
#     )

# def extract_docx_text(path):
#     return "\n".join(
#         p.text 
#         for p in docx.Document(path).paragraphs 
#         if p.text.strip()
#     )

# if __name__ == "__main__":
#     process_files()
import os
import shutil
import io
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
from pathlib import Path
import pytesseract
from PIL import Image
from pdf2image import convert_from_path
pytesseract.pytesseract.tesseract_cmd = '/opt/homebrew/bin/tesseract'

def process_files(input_dir="docs", output_dir="text_files"):
    """Process all supported files (PDF, PPTX, DOCX, TXT, Images)."""
    Path(output_dir).mkdir(exist_ok=True)
    
    for filename in os.listdir(input_dir):
        filepath = os.path.join(input_dir, filename)
        output_path = os.path.join(output_dir, f"{Path(filename).stem}.txt")
        
        try:
            if filename.lower().endswith('.pdf'):
                text = extract_pdf_text(filepath)
            elif filename.lower().endswith('.pptx'):
                text = extract_pptx_text(filepath)
            elif filename.lower().endswith('.docx'):
                text = extract_docx_text(filepath)
            elif filename.lower().endswith('.txt'):
                shutil.copy2(filepath, output_path)
                print(f"Copied: {filename} → {output_path}")
                continue
            elif filename.lower().split('.')[-1] in ('png', 'jpg', 'jpeg', 'bmp', 'tiff'):
                text = extract_image_text(filepath)
            else:
                continue
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            print(f"Processed: {filename} → {output_path}")
            
        except Exception as e:
            print(f"Failed {filename}: {str(e)}")

def extract_pdf_text(path):
    """Extract text from PDF (plain text + images via OCR)."""
    text = []
    with open(path, 'rb') as f:
        reader = PdfReader(f)
        for i, page in enumerate(reader.pages):
            # Extract plain text
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
            
            # Extract text from images in PDF
            if '/XObject' in page['/Resources']:
                x_objects = page['/Resources']['/XObject'].get_object()
                for obj in x_objects:
                    if x_objects[obj]['/Subtype'] == '/Image':
                        images = convert_from_path(
                            path, 
                            first_page=i+1,
                            last_page=i+1,
                            dpi=300
                        )
                        for img in images:
                            text.append(pytesseract.image_to_string(img))
    return "\n".join(text)

def extract_pptx_text(path):
    """Extract text from PPTX (slides + images via OCR)."""
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            # Extract plain text
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text.strip())
            
            # Extract text from images
            if shape.shape_type == 13:  # 13 = picture type
                img_bytes = shape.image.blob
                with Image.open(io.BytesIO(img_bytes)) as img:
                    img_text = pytesseract.image_to_string(img)
                    if img_text.strip():
                        text.append(img_text)
    return "\n".join(text)

def extract_docx_text(path):
    """Extract text from DOCX (plain text only)."""
    return "\n".join(
        p.text 
        for p in docx.Document(path).paragraphs 
        if p.text.strip()
    )

def extract_image_text(path):
    """Extract text from standalone images."""
    return pytesseract.image_to_string(Image.open(path))

if __name__ == "__main__":
    process_files()